"""Render a project's synthesis as a customer-ready .pptx deck.

Prefers a deck-outline artifact when present; otherwise falls back to a
generic one-slide-per-artifact layout grouped by category.
"""

from __future__ import annotations

from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from app.synthesis.exporters.common import (
    category_label,
    group_by_category,
    list_items,
    short,
)
from app.synthesis.models import Artifact

_TITLE_LAYOUT = 0
_TITLE_AND_CONTENT_LAYOUT = 1
_SECTION_HEADER_LAYOUT = 2


def export_pptx(project: dict, artifacts: list[Artifact]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    customer = str(project.get("customer") or project.get("name") or "Customer").strip()
    project_name = str(project.get("name") or project.get("slug") or "Project").strip()

    _add_title_slide(prs, customer, project_name)

    outline = _find_deck_outline(artifacts)
    if outline:
        _render_outline(prs, outline)
    else:
        _render_default(prs, artifacts)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_title_slide(prs: Presentation, customer: str, project_name: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_LAYOUT])
    if slide.shapes.title:
        slide.shapes.title.text = f"{customer}"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = project_name


def _find_deck_outline(artifacts: list[Artifact]) -> Artifact | None:
    for a in artifacts:
        if a.type_id == "deck-outline":
            return a
    return None


def _render_outline(prs: Presentation, outline: Artifact) -> None:
    slides = list_items(outline.body.get("slides"))
    for slide_def in slides:
        if not isinstance(slide_def, dict):
            continue
        title = str(slide_def.get("title") or "").strip() or "Slide"
        bullets = _bullets_from(slide_def)
        _add_content_slide(prs, title, bullets)


def _render_default(prs: Presentation, artifacts: list[Artifact]) -> None:
    grouped = group_by_category(artifacts)
    for category, group in grouped.items():
        if not group:
            continue
        _add_section_slide(prs, category_label(category))
        for artifact in group:
            bullets: list[str] = []
            if artifact.summary:
                bullets.append(short(artifact.summary, 280))
            for value in artifact.body.values():
                items = list_items(value)
                for item in items[:6]:
                    bullets.append(short(item, 200))
                if len(bullets) >= 7:
                    break
            _add_content_slide(prs, artifact.title or artifact.type_id, bullets[:7])


def _bullets_from(slide_def: dict[str, Any]) -> list[str]:
    bullets: list[str] = []
    key_point = slide_def.get("key_point")
    if key_point:
        bullets.append(short(key_point, 240))
    for item in list_items(slide_def.get("supporting_bullets")):
        bullets.append(short(item, 200))
    return bullets[:7]


def _add_section_slide(prs: Presentation, title: str) -> None:
    layout_idx = _SECTION_HEADER_LAYOUT
    if layout_idx >= len(prs.slide_layouts):
        layout_idx = _TITLE_LAYOUT
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    if slide.shapes.title:
        slide.shapes.title.text = title


def _add_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_AND_CONTENT_LAYOUT])
    if slide.shapes.title:
        slide.shapes.title.text = title

    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:
            body_ph = ph
            break

    if body_ph is None or not bullets:
        return

    tf = body_ph.text_frame
    tf.word_wrap = True
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(16)
