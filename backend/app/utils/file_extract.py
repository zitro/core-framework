"""Convert uploaded binary documents (PDF/DOCX/PPTX/XLSX) to markdown text.

Used by the file-upload ingest endpoint to turn arbitrary office files into
something the LLM classifier can route into the markdown engagement repo.

Each extractor returns a single markdown string. Failures raise
`UnsupportedFileType` or `ExtractionError`.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


class UnsupportedFileTypeError(ValueError):
    """Raised when the file extension is not supported."""


class ExtractionError(RuntimeError):
    """Raised when a supported file fails to extract."""


# Map of lowercase extension (with dot) -> extractor name for logging.
SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".md",
    ".markdown",
    ".txt",
}


def _extract_pdf(data: bytes) -> str:
    import pymupdf  # lazy import: keep startup fast

    try:
        doc = pymupdf.open(stream=data, filetype="pdf")
    except Exception as exc:  # pragma: no cover - corrupted file path
        raise ExtractionError(f"failed to open PDF: {exc}") from exc

    parts: list[str] = []
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            parts.append(f"## Page {page_num}\n\n{text}")
    doc.close()
    return "\n\n".join(parts).strip()


def _extract_docx(data: bytes) -> str:
    from docx import Document

    try:
        doc = Document(io.BytesIO(data))
    except Exception as exc:
        raise ExtractionError(f"failed to open DOCX: {exc}") from exc

    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading 1"):
            parts.append(f"# {text}")
        elif style.startswith("heading 2"):
            parts.append(f"## {text}")
        elif style.startswith("heading 3"):
            parts.append(f"### {text}")
        else:
            parts.append(text)

    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            header = rows[0]
            sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
            parts.append("\n".join([header, sep, *rows[1:]]))

    return "\n\n".join(parts).strip()


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise ExtractionError(f"failed to open PPTX: {exc}") from exc

    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = [f"## Slide {idx}"]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_lines.append(f"- {text}" if para.level > 0 else text)
        if len(slide_lines) > 1:
            parts.append("\n".join(slide_lines))
    return "\n\n".join(parts).strip()


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    try:
        wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    except Exception as exc:
        raise ExtractionError(f"failed to open XLSX: {exc}") from exc

    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        # Trim trailing all-None rows
        while rows and all(c is None for c in rows[-1]):
            rows.pop()
        if not rows:
            continue
        widths = [
            max(
                1,
                max(
                    (len(str(row[i])) if i < len(row) and row[i] is not None else 0) for row in rows
                ),
            )
            for i in range(max(len(r) for r in rows))
        ]
        parts.append(f"## {sheet_name}")
        header = rows[0]
        header_cells = [str(c) if c is not None else "" for c in header]
        sep_cells = ["---"] * len(widths)
        parts.append("| " + " | ".join(header_cells) + " |")
        parts.append("| " + " | ".join(sep_cells) + " |")
        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            parts.append("| " + " | ".join(cells) + " |")
    wb.close()
    return "\n\n".join(parts).strip()


def _extract_text(data: bytes) -> str:
    return data.decode("utf-8", errors="replace").strip()


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".md": _extract_text,
    ".markdown": _extract_text,
    ".txt": _extract_text,
}


def extract_to_markdown(filename: str, data: bytes) -> str:
    """Convert an uploaded file's bytes into a markdown string.

    Raises:
        UnsupportedFileTypeError: If the extension is not in `SUPPORTED_EXTENSIONS`.
        ExtractionError: If the file is supported but extraction fails.
    """
    ext = Path(filename).suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise UnsupportedFileTypeError(
            f"Unsupported file type: {ext or '(none)'}. Supported: {sorted(SUPPORTED_EXTENSIONS)}"
        )
    text = extractor(data)
    if not text:
        raise ExtractionError(f"no text extracted from {filename}")
    return text
