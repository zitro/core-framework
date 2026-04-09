"""Read local document files from a directory for AI context."""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# Text formats — read directly as UTF-8
TEXT_EXTENSIONS = {
    ".md", ".txt", ".csv", ".json", ".yaml", ".yml",
    ".rst", ".log", ".html", ".xml", ".toml", ".ini",
}

# Binary formats — need dedicated parsers
BINARY_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".xlsx",
}

ALLOWED_EXTENSIONS = TEXT_EXTENSIONS | BINARY_EXTENSIONS
MAX_FILE_SIZE = 200_000  # 200 KB extracted text per file
MAX_TOTAL_SIZE = 500_000  # 500 KB combined


def _read_pdf(path: Path) -> str:
    import fitz  # pymupdf

    text_parts: list[str] = []
    with fitz.open(path) as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts)


def _read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(max_row=200, values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


BINARY_READERS = {
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".pptx": _read_pptx,
    ".xlsx": _read_xlsx,
}


def validate_docs_path(path_str: str) -> Path:
    """Resolve and validate a docs directory path.

    Returns the resolved Path. Raises ValueError if invalid.
    """
    if not path_str or not path_str.strip():
        raise ValueError("Empty docs path")
    resolved = Path(path_str).expanduser().resolve()
    if not resolved.exists():
        raise ValueError(f"Path does not exist: {resolved}")
    if not resolved.is_dir():
        raise ValueError(f"Path is not a directory: {resolved}")
    return resolved


def scan_docs(path_str: str) -> list[dict]:
    """List readable document files in the directory.

    Returns list of {name, size, extension} dicts.
    """
    docs_dir = validate_docs_path(path_str)
    files: list[dict] = []
    for f in sorted(docs_dir.rglob("*")):
        if not f.is_file():
            continue
        if f.suffix.lower() not in ALLOWED_EXTENSIONS:
            continue
        # Skip hidden files and common non-doc directories
        parts = f.relative_to(docs_dir).parts
        if any(p.startswith(".") or p in ("node_modules", "__pycache__") for p in parts):
            continue
        files.append({
            "name": str(f.relative_to(docs_dir)),
            "size": f.stat().st_size,
            "extension": f.suffix.lower(),
        })
    return files


def read_docs_content(path_str: str) -> str:
    """Read all eligible documents and return combined text.

    Respects per-file and total size caps.
    """
    docs_dir = validate_docs_path(path_str)
    files = scan_docs(path_str)
    parts: list[str] = []
    total = 0

    for info in files:
        if total >= MAX_TOTAL_SIZE:
            parts.append(f"[Truncated — {MAX_TOTAL_SIZE // 1000}KB limit reached]")
            break
        fpath = docs_dir / info["name"]
        ext = info["extension"]
        try:
            if ext in BINARY_READERS:
                raw = BINARY_READERS[ext](fpath)
            else:
                raw = fpath.read_text(encoding="utf-8", errors="replace")
        except Exception:
            logger.debug("Skipping unreadable file: %s", fpath)
            continue
        if len(raw) > MAX_FILE_SIZE:
            raw = raw[:MAX_FILE_SIZE] + "\n[...file truncated]"
        parts.append(f"--- {info['name']} ---\n{raw}")
        total += len(raw)

    return "\n\n".join(parts)
