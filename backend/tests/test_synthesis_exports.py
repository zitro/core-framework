"""Phase 6J router tests for POST /api/synthesis/{project_id}/export/{fmt}."""

from __future__ import annotations

from io import BytesIO

import pytest
from httpx import AsyncClient

from app.providers.storage import get_storage_provider

pytestmark = pytest.mark.asyncio


async def _seed_project(project_id: str, name: str = "Claims rework") -> None:
    storage = get_storage_provider()
    await storage.create(
        "engagements",
        {
            "id": project_id,
            "name": name,
            "customer": "Acme",
            "status": "proposed",
            "owners": [],
        },
    )


async def _seed_artifact(
    project_id: str,
    artifact_id: str,
    *,
    type_id: str = "problem-statement",
    category: str = "why",
    title: str = "Claims pain",
    body: dict | None = None,
) -> None:
    storage = get_storage_provider()
    await storage.create(
        "artifacts",
        {
            "id": artifact_id,
            "project_id": project_id,
            "type_id": type_id,
            "category": category,
            "title": title,
            "summary": "Claims agents juggle three tools.",
            "body": body or {"statement": "Agents juggle 3 tools."},
            "citations": [],
            "status": "draft",
            "version": 1,
        },
    )


# ── docx ───────────────────────────────────────────────────────────────


async def test_export_docx_returns_zip_payload(client: AsyncClient) -> None:
    await _seed_project("proj-x-1")
    await _seed_artifact("proj-x-1", "art-x-1")
    resp = await client.post("/api/synthesis/proj-x-1/export/docx")
    assert resp.status_code == 200
    assert resp.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
    assert "claims-rework-synthesis.docx" in resp.headers["content-disposition"]
    # .docx files are ZIP containers — first 2 bytes are "PK".
    assert resp.content[:2] == b"PK"


async def test_export_docx_with_no_artifacts_still_succeeds(client: AsyncClient) -> None:
    await _seed_project("proj-x-empty")
    resp = await client.post("/api/synthesis/proj-x-empty/export/docx")
    assert resp.status_code == 200
    assert resp.content[:2] == b"PK"


# ── pptx ───────────────────────────────────────────────────────────────


async def test_export_pptx_returns_zip_payload(client: AsyncClient) -> None:
    await _seed_project("proj-p-1")
    await _seed_artifact("proj-p-1", "art-p-1")
    resp = await client.post("/api/synthesis/proj-p-1/export/pptx")
    assert resp.status_code == 200
    assert resp.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert "claims-rework-synthesis.pptx" in resp.headers["content-disposition"]
    assert resp.content[:2] == b"PK"


async def test_export_pptx_uses_deck_outline_when_present(client: AsyncClient) -> None:
    """A deck-outline artifact should drive slide bullets — verify the
    payload is a valid pptx and contains at least the title slide."""
    await _seed_project("proj-p-outline")
    await _seed_artifact(
        "proj-p-outline",
        "art-outline",
        type_id="deck-outline",
        category="what",
        title="Customer readout",
        body={
            "slides": [
                {
                    "title": "Why we're here",
                    "key_point": "Claims cycle time is the constraint",
                    "supporting_bullets": ["Agents juggle 3 tools", "Hand-offs lose context"],
                },
                {
                    "title": "What we'll do",
                    "key_point": "Consolidate intake",
                    "supporting_bullets": ["Single queue", "Unified UI"],
                },
            ]
        },
    )
    resp = await client.post("/api/synthesis/proj-p-outline/export/pptx")
    assert resp.status_code == 200
    assert resp.content[:2] == b"PK"

    # The pptx is a real Office Open XML zip — round-trip through python-pptx
    # so we know the file is well-formed and contains our outline slides.
    from pptx import Presentation

    prs = Presentation(BytesIO(resp.content))
    titles = [
        s.shapes.title.text for s in prs.slides if s.shapes.title and s.shapes.title.text
    ]
    assert "Why we're here" in titles
    assert "What we'll do" in titles


# ── errors ─────────────────────────────────────────────────────────────


async def test_export_404_when_project_missing(client: AsyncClient) -> None:
    resp = await client.post("/api/synthesis/missing-id/export/docx")
    assert resp.status_code == 404


async def test_export_400_when_format_unsupported(client: AsyncClient) -> None:
    await _seed_project("proj-bad-fmt")
    resp = await client.post("/api/synthesis/proj-bad-fmt/export/pdf")
    assert resp.status_code == 400
